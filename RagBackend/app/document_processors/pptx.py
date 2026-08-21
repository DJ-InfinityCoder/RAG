"""
PPTX extraction: slide-by-slide with titles, tables, and speaker notes.
"""

from typing import List
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.config import get_logger
from app.document_processors.base import format_table_as_markdown

logger = get_logger("askdoc.processors.pptx")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_pptx_slides(pptx_path: str, filename: str, chunk_size: int = 1800, chunk_overlap: int = 300) -> List[Document]:
    """
    Extracts PPTX content slide-by-slide preserving structure:
    1. Extracts slide title, body shapes, and inlined table shapes (via shape.has_table).
    2. Appends speaker notes from slide.notes_slide if present.
    3. Prepends clear slide weighting (e.g. 'Slide {slide_num}: {title}').
    4. Creates 1 chunk per slide by default, only splitting if content exceeds chunk_size.
    5. Attaches slide_number, page, and has_notes metadata.
    """
    if not Presentation:
        logger.warning("python-pptx not available, cannot parse PPTX")
        return []

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error(f"Error opening PPTX file {filename}: {e}", exc_info=True)
        return []

    documents = []
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n\n", "\n\n", "\n", ". ", "? ", "! ", " ", ""]
    )

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        
        # 1. Identify slide title
        slide_title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()
        except Exception:
            slide_title = ""

        # 2. Extract shape contents (text and tables)
        body_elements = []
        for shape in slide.shapes:
            # Skip title shape if already extracted
            if slide.shapes.title and shape == slide.shapes.title:
                continue

            # Check for Table Shape
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    if any(row_cells):
                        table_data.append(row_cells)
                if table_data:
                    md_table = format_table_as_markdown(table_data)
                    if md_table.strip():
                        body_elements.append(md_table)
            # Check for Text Frame
            elif shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    body_elements.append(txt)
            elif hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    body_elements.append(txt)

        # 3. Extract Speaker Notes
        notes_text = ""
        has_notes = False
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    has_notes = True
        except Exception:
            pass

        # 4. Construct Slide Content with strong title prefix
        header_line = f"# Slide {slide_num}" + (f": {slide_title}" if slide_title else "")
        content_parts = [header_line]
        if body_elements:
            content_parts.append("\n\n".join(body_elements))
        if notes_text:
            content_parts.append(f"**Speaker Notes:**\n{notes_text}")

        slide_full_text = "\n\n".join(content_parts).strip()
        if not slide_full_text:
            continue

        slide_meta = {
            "source": filename,
            "title": filename,
            "slide_number": slide_num,
            "page": slide_num,
            "page_number": slide_num,
            "has_notes": has_notes,
            "slide_title": slide_title
        }

        # 5. One chunk per slide by default; split only if too large
        if len(slide_full_text) <= chunk_size:
            documents.append(Document(page_content=slide_full_text, metadata=dict(slide_meta)))
        else:
            split_texts = splitter.split_text(slide_full_text)
            for chunk_part_idx, part in enumerate(split_texts):
                meta = dict(slide_meta)
                meta["part_index"] = chunk_part_idx
                documents.append(Document(page_content=part, metadata=meta))

    return documents
