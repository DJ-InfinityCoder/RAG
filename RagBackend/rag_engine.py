import os
import time
from typing import List, Optional
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_pinecone import PineconeVectorStore
from pinecone import Pinecone
from langchain_core.documents import Document
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.messages import HumanMessage, AIMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, UnstructuredPowerPointLoader
from langchain_core.embeddings import Embeddings
from functools import lru_cache
import tempfile
import pandas as pd
import json
from flashrank import Ranker, RerankRequest
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

class PineconeInferenceEmbeddings(Embeddings):
    def __init__(self, api_key: str, model: str = "llama-text-embed-v2"):
        self.pc = Pinecone(api_key=api_key)
        self.model = model

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        try:
            batch_size = 90
            all_embeddings = []
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i + batch_size]
                response = self.pc.inference.embed(
                    model=self.model,
                    inputs=batch,
                    parameters={"input_type": "passage", "truncate": "END"}
                )
                all_embeddings.extend([r['values'] for r in response])
            return all_embeddings
        except Exception as e:
            print(f"Error embedding documents: {e}")
            raise e

    @lru_cache(maxsize=1000)
    def embed_query(self, text: str) -> List[float]:
        try:
            response = self.pc.inference.embed(
                model=self.model,
                inputs=[text],
                parameters={"input_type": "query", "truncate": "END"}
            )
            return response[0]['values']
        except Exception as e:
            print(f"Error embedding query: {e}")
            raise e

class RAGEngine:
    def __init__(self):
        GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
        PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
        PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME", "djrag")

        # Use Pinecone Inference for embeddings (1024 dimensions)
        self.embeddings = PineconeInferenceEmbeddings(api_key=PINECONE_API_KEY, model="llama-text-embed-v2")
        
        # Use Gemini for Chat
        self.llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", google_api_key=GOOGLE_API_KEY)
        
        # Initialize Pinecone
        self.pc = Pinecone(api_key=PINECONE_API_KEY)
        
        # Connect to existing index
        self.index = self.pc.Index(PINECONE_INDEX_NAME)
        self.vectorstore = PineconeVectorStore(embedding=self.embeddings, index=self.index)
        
        # Initialize FlashRank
        self.ranker = Ranker()

        # Rephrasing Prompt
        self.rephrase_prompt = ChatPromptTemplate.from_messages([
            ("system", """Given the following conversation and a follow-up question, rephrase the follow-up question to be a standalone question.
Chat History:
{chat_history}
Follow Up Input: {question}
Standalone question:"""),
        ])

        # Main QA Prompt
        self.prompt = ChatPromptTemplate.from_messages([
            ("system", """You are a professional and knowledgeable AI assistant. Use the following context to answer the question.
- Provide a comprehensive and structured answer.
- Use Markdown formatting (headers, bullet points, bold text) to enhance readability.
- If the answer is not in the context, politely state that you don't have that information.
- Always provide inline citations for your answer using the format [1], [2], etc.
- These citations must correspond to the source documents provided in the context.

Context:
{context}"""),
            MessagesPlaceholder(variable_name="chat_history"),
            ("user", "{question}"),
        ])

    def chat(self, question: str, session_id: str = None, chat_history: List = []):
        # 0. Rephrase question if there is history
        search_query = question
        if chat_history:
            rephrase_chain = self.rephrase_prompt | self.llm
            try:
                # Convert LangChain messages to string for rephrasing context
                history_str = "\n".join([f"{m.type}: {m.content}" for m in chat_history])
                response = rephrase_chain.invoke({"chat_history": history_str, "question": question})
                search_query = response.content
                print(f"Rephrased query: {search_query}")
            except Exception as e:
                print(f"Error rephrasing query: {e}")

        # 1. Retrieval
        filter_dict = {}
        if session_id:
            filter_dict = {"session_id": session_id}
            
        docs = self.vectorstore.similarity_search(search_query, k=10, filter=filter_dict)
        
        # 2. Reranking
        passages = [
            {"id": str(i), "text": doc.page_content, "meta": doc.metadata} 
            for i, doc in enumerate(docs)
        ]
        
        if passages:
            rerank_request = RerankRequest(query=search_query, passages=passages)
            results = self.ranker.rerank(rerank_request)
            top_results = results[:5]
            
            reranked_docs = []
            for res in top_results:
                doc = Document(page_content=res['text'], metadata=res['meta'])
                reranked_docs.append(doc)
        else:
            reranked_docs = []

        # 3. Generation
        context_text = ""
        sources = []
        for i, doc in enumerate(reranked_docs):
            index = i + 1
            context_text += f"Source [{index}]:\n{doc.page_content}\n\n"
            sources.append({
                "id": index,
                "title": doc.metadata.get("source", "Unknown"),
                "content": doc.page_content,
                "metadata": doc.metadata
            })

        messages = self.prompt.invoke({
            "question": question, # Use original question for generation to maintain tone
            "context": context_text,
            "chat_history": chat_history
        })
        start_time = time.time()
        response = self.llm.invoke(messages)
        end_time = time.time()
        duration = end_time - start_time
        
        # Calculate metrics
        input_tokens = 0
        output_tokens = 0
        if response.usage_metadata:
            input_tokens = response.usage_metadata.get("input_tokens", 0)
            output_tokens = response.usage_metadata.get("output_tokens", 0)
            
        # Cost estimation for Gemini 1.5 Flash (approximate)
        # Input: $0.075 / 1M tokens
        # Output: $0.30 / 1M tokens
        input_cost = (input_tokens / 1_000_000) * 0.075
        output_cost = (output_tokens / 1_000_000) * 0.30
        total_cost = input_cost + output_cost
        
        metrics = {
            "time": round(duration, 2),
            "input_tokens": input_tokens,
            "output_tokens": output_tokens,
            "total_tokens": input_tokens + output_tokens,
            "cost": round(total_cost, 6)
        }
        
        return {"answer": response.content, "sources": sources, "metrics": metrics}

    def ingest_text(self, text: str, title: str, session_id: str = None):
        if not text.strip():
            return 0
            
        documents = [Document(page_content=text, metadata={"source": title})]
        
        # Chunking strategy: ~1000 tokens (approx 4000 chars) with 10% overlap
        splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=400)
        docs = splitter.split_documents(documents)
        
        if session_id:
            for i, doc in enumerate(docs):
                doc.metadata["session_id"] = session_id
                doc.metadata["title"] = title
                doc.metadata["chunk_index"] = i
                doc.metadata["total_chunks"] = len(docs)
                
        self.vectorstore.add_documents(documents=docs)
        return len(docs)

    async def process_file(self, file_content: bytes, filename: str, session_id: str = None):
        # Create a temp file to save the uploaded content
        suffix = os.path.splitext(filename)[1]
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name

        try:
            documents = []
            filename_lower = filename.lower()
            
            if filename_lower.endswith(".pdf"):
                loader = PyPDFLoader(tmp_path)
                documents = loader.load()
            elif filename_lower.endswith(".docx"):
                loader = Docx2txtLoader(tmp_path)
                documents = loader.load()
            elif filename_lower.endswith(".pptx"):
                # Simple PPTX extraction using python-pptx
                if Presentation:
                    prs = Presentation(tmp_path)
                    text_content = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content.append(shape.text)
                    full_text = "\n".join(text_content)
                    documents = [Document(page_content=full_text, metadata={"source": filename})]
            elif filename_lower.endswith(".txt"):
                with open(tmp_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                documents = [Document(page_content=text, metadata={"source": filename})]
            elif filename_lower.endswith(".xlsx"):
                df = pd.read_excel(tmp_path)
                for _, row in df.iterrows():
                    content = "\n".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    documents.append(Document(page_content=content, metadata={"source": filename, "row": _}))
            elif filename_lower.endswith(".csv"):
                df = pd.read_csv(tmp_path)
                for _, row in df.iterrows():
                    content = "\n".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    documents.append(Document(page_content=content, metadata={"source": filename, "row": _}))
            
            if documents:
                # Chunking strategy: ~1000 tokens (approx 4000 chars) with 10% overlap
                splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=400)
                docs = splitter.split_documents(documents)
                
                if session_id:
                    for i, doc in enumerate(docs):
                        doc.metadata["session_id"] = session_id
                        doc.metadata["title"] = filename
                        doc.metadata["chunk_index"] = i
                        doc.metadata["total_chunks"] = len(docs)
                        
                self.vectorstore.add_documents(documents=docs)
                return len(docs)
            return 0
        except Exception as e:
            print(f"Error processing file {filename}: {e}")
            raise e
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    def delete_vectors(self, session_id: str):
        try:
            self.index.delete(filter={"session_id": session_id})
            return True
        except Exception as e:
            print(f"Error deleting vectors: {e}")
            return False
